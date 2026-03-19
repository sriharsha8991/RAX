"""
RAX — Resume Analysis eXpert | Stage 2 PPT Generator
Generates a clean 6-slide PowerPoint matching the requirement spec:
  Slide 1: Title, Team, Target Market, Value Propositions
  Slide 2: Competitive Landscape & Winning Feature #1
  Slide 3: Winning Features #2 & #3
  Slide 4: Winning Feature #4 & Workflows
  Slide 5: System Architecture
  Slide 6: Conclusion & Implementation Plan
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color Palette ──────────────────────────────────────────────────
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG     = RGBColor(0x1A, 0x1A, 0x2E)  # Deep navy
ACCENT_BLUE = RGBColor(0x00, 0x7A, 0xCC)  # Primary blue
ACCENT_TEAL = RGBColor(0x00, 0xB4, 0xD8)  # Teal accent
ACCENT_GREEN= RGBColor(0x2E, 0xCC, 0x71)  # Green
ACCENT_RED  = RGBColor(0xE7, 0x4C, 0x3C)  # Red for competitors
LIGHT_GRAY  = RGBColor(0xE8, 0xE8, 0xE8)
MED_GRAY    = RGBColor(0x99, 0x99, 0x99)
DARK_TEXT    = RGBColor(0x2C, 0x2C, 0x2C)
BODY_TEXT    = RGBColor(0x44, 0x44, 0x44)
ORANGE      = RGBColor(0xE6, 0x7E, 0x22)
PURPLE      = RGBColor(0x9B, 0x59, 0xB6)
PINK        = RGBColor(0xE9, 0x1E, 0x63)
SLIDE_BG    = RGBColor(0xFF, 0xFF, 0xFF)
HEADER_BG   = RGBColor(0x00, 0x52, 0x8A)  # Darker blue for headers

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def set_text(shape, text, font_size=14, bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_paragraph(tf, text, font_size=14, bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT, space_before=0, space_after=0, font_name="Calibri"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    return p


def add_text_box(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    return txBox


def add_bullet_list(tf, items, font_size=13, color=BODY_TEXT, bold_first_part=True):
    """Add bulleted items. If bold_first_part, text before ' — ' is bolded."""
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        p.space_after = Pt(2)
        p.level = 0

        if bold_first_part and " — " in item:
            parts = item.split(" — ", 1)
            run1 = p.add_run()
            run1.text = "• " + parts[0] + " — "
            run1.font.size = Pt(font_size)
            run1.font.bold = True
            run1.font.color.rgb = color
            run1.font.name = "Calibri"
            run2 = p.add_run()
            run2.text = parts[1]
            run2.font.size = Pt(font_size)
            run2.font.bold = False
            run2.font.color.rgb = color
            run2.font.name = "Calibri"
        else:
            run = p.add_run()
            run.text = "• " + item
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.name = "Calibri"


def add_table(slide, left, top, width, height, rows, cols, data, col_widths=None, header_color=HEADER_BG):
    """Add a styled table. data = list of lists. First row = header."""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c] if r < len(data) and c < len(data[r]) else ""

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.name = "Calibri"
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.font.color.rgb = DARK_TEXT
                    paragraph.alignment = PP_ALIGN.LEFT

            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF8, 0xFC)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

    return table


# ══════════════════════════════════════════════════════════════════
#  SLIDE 1 — Title, Team, Target Market, Value Propositions
# ══════════════════════════════════════════════════════════════════
def slide_1_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, WHITE)

    # Top banner
    add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(2.6), DARK_BG)

    # Title
    tb = add_text_box(slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.7))
    set_text(tb.shape, "RAX — Resume Analysis eXpert", font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)

    # Subtitle
    tb2 = add_text_box(slide, Inches(0.6), Inches(1.05), Inches(12), Inches(0.5))
    set_text(tb2.shape, "AI-Powered Hiring  •  Explainable  •  Bias-Free  •  Real-Time  •  Distributed", font_size=18, color=ACCENT_TEAL, alignment=PP_ALIGN.LEFT)

    # Team
    tb3 = add_text_box(slide, Inches(0.6), Inches(1.7), Inches(12), Inches(0.5))
    set_text(tb3.shape, "Team:  Person 1 (Backend Core)  ·  Person 2 (AI Pipeline)  ·  Person 3 (Frontend)", font_size=14, color=RGBColor(0xBB, 0xBB, 0xBB), alignment=PP_ALIGN.LEFT)

    # ── Target Market section ──
    tb4 = add_text_box(slide, Inches(0.6), Inches(2.9), Inches(5.8), Inches(0.45))
    set_text(tb4.shape, "TARGET MARKET", font_size=16, bold=True, color=ACCENT_BLUE)

    market_items = [
        "Enterprise HR teams — handling 100+ resumes per open role",
        "Staffing agencies & RPO firms — thousands of resumes across clients",
        "Remote-first & distributed teams — need real-time collaborative screening",
        "Compliance-driven orgs — DEI mandates, EEOC regulations",
    ]
    tb5 = add_text_box(slide, Inches(0.6), Inches(3.35), Inches(5.8), Inches(2.5))
    tf = tb5.text_frame
    tf.word_wrap = True
    add_bullet_list(tf, market_items, font_size=12, color=BODY_TEXT)

    # ── Value Propositions section ──
    tb6 = add_text_box(slide, Inches(6.8), Inches(2.9), Inches(6.2), Inches(0.45))
    set_text(tb6.shape, "VALUE PROPOSITIONS", font_size=16, bold=True, color=ACCENT_BLUE)

    vp_items = [
        "Semantic understanding — replaces keyword matching",
        "Explainable scoring — strengths, gaps, reasoning for every candidate",
        "Bias-aware screening — names, gender, institutions anonymized before evaluation",
        "Hybrid Graph + Vector — Neo4j structural reasoning + Qdrant semantic discovery",
        "Real-time pipeline — WebSocket-driven live processing visibility",
        "AI candidate feedback — constructive rejection feedback for employer brand",
        "Free-tier cloud — 100% deployed on free-tier services",
    ]
    tb7 = add_text_box(slide, Inches(6.8), Inches(3.35), Inches(6.2), Inches(3.8))
    tf7 = tb7.text_frame
    tf7.word_wrap = True
    add_bullet_list(tf7, vp_items, font_size=12, color=BODY_TEXT)

    # Bottom accent line
    add_shape_rect(slide, Inches(0), Inches(7.25), SLIDE_WIDTH, Inches(0.05), ACCENT_BLUE)

    # Slide number
    tb_num = add_text_box(slide, Inches(12.4), Inches(7.05), Inches(0.8), Inches(0.35))
    set_text(tb_num.shape, "1 / 6", font_size=10, color=MED_GRAY, alignment=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════
#  SLIDE 2 — Competitive Landscape & Winning Feature #1
# ══════════════════════════════════════════════════════════════════
def slide_2_competitive(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # Header bar
    add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.95), DARK_BG)
    tb = add_text_box(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.6))
    set_text(tb.shape, "Competitive Landscape & Winning Feature #1: Hybrid Graph + Vector Matching", font_size=24, bold=True, color=WHITE)

    # ── Left: Competitor Table ──
    tb_label = add_text_box(slide, Inches(0.4), Inches(1.15), Inches(5), Inches(0.4))
    set_text(tb_label.shape, "COMPETITOR GAPS", font_size=14, bold=True, color=ACCENT_RED)

    comp_data = [
        ["Competitor", "Strengths", "Key Limitation"],
        ["Greenhouse", "Robust ATS workflow, 400+ integrations", "Keyword-only filtering; no AI explanation"],
        ["Lever", "CRM + ATS hybrid, collaboration", "Scoring is manual/opaque; no anonymization"],
        ["HireVue", "AI-powered video assessment", "Video-only; resume screening is secondary"],
        ["Pymetrics", "Neuroscience-based games", "Requires candidate participation; not resume-centric"],
        ["Ideal (Ceridian)", "AI screening & shortlisting", "Black-box proprietary; zero explainability"],
        ["Textio", "JD language optimization", "Only JD side; doesn't score candidates"],
    ]
    add_table(slide, Inches(0.4), Inches(1.55), Inches(5.8), Inches(3.0), 7, 3, comp_data,
              col_widths=[Inches(1.3), Inches(2.1), Inches(2.4)], header_color=ACCENT_RED)

    # gap callout
    tb_gap = add_text_box(slide, Inches(0.4), Inches(4.7), Inches(5.8), Inches(0.65))
    tf_gap = tb_gap.text_frame
    tf_gap.word_wrap = True
    p = tf_gap.paragraphs[0]
    run = p.add_run()
    run.text = "Shared gap: "
    run.font.bold = True
    run.font.size = Pt(12)
    run.font.color.rgb = ACCENT_RED
    run.font.name = "Calibri"
    run2 = p.add_run()
    run2.text = "Opaque scoring + no structural reasoning + bias anonymization is an afterthought."
    run2.font.size = Pt(12)
    run2.font.color.rgb = BODY_TEXT
    run2.font.name = "Calibri"

    # ── Right: Winning Feature #1 ──
    tb_wf = add_text_box(slide, Inches(6.6), Inches(1.15), Inches(6.4), Inches(0.4))
    set_text(tb_wf.shape, "WINNING FEATURE #1: HYBRID GRAPH + VECTOR", font_size=14, bold=True, color=ACCENT_BLUE)

    hybrid_data = [
        ["Capability", "Vector-Only (Others)", "RAX (Hybrid)"],
        ["Find similar resumes", "✅ Cosine similarity", "✅ Qdrant embeddings"],
        ["Explain WHY it matches", "❌ Just a score", "✅ Neo4j graph paths"],
        ["Adjacent skill discovery", "Partial", "✅ Embeddings + IS_SIMILAR_TO edges"],
        ["Multi-hop reasoning", "❌ Impossible", "✅ Graph traversal"],
        ["Skill taxonomy", "❌ Flat embeddings", "✅ Hierarchical graph"],
        ["Bias audit trail", "❌ Can't verify", "✅ Queryable graph"],
        ["Self-improving", "❌ Static model", "✅ Auto-enriching graph"],
    ]
    add_table(slide, Inches(6.6), Inches(1.55), Inches(6.4), Inches(3.4), 8, 3, hybrid_data,
              col_widths=[Inches(1.9), Inches(2.0), Inches(2.5)], header_color=ACCENT_BLUE)

    # Score fusion formula box
    formula_box = add_rounded_rect(slide, Inches(6.6), Inches(5.15), Inches(6.4), Inches(1.1), RGBColor(0xEE, 0xF5, 0xFF), ACCENT_BLUE)
    tf_f = formula_box.text_frame
    tf_f.word_wrap = True
    p = tf_f.paragraphs[0]
    run = p.add_run()
    run.text = "Score Fusion Formula:"
    run.font.bold = True
    run.font.size = Pt(12)
    run.font.color.rgb = ACCENT_BLUE
    run.font.name = "Calibri"
    p2 = tf_f.add_paragraph()
    p2.space_before = Pt(6)
    run2 = p2.add_run()
    run2.text = "hybrid = 0.30 × semantic (Qdrant) + 0.50 × structural (Neo4j) + 0.15 × experience + 0.05 × education"
    run2.font.size = Pt(11)
    run2.font.color.rgb = DARK_TEXT
    run2.font.name = "Consolas"
    p3 = tf_f.add_paragraph()
    p3.space_before = Pt(4)
    run3 = p3.add_run()
    run3.text = "Weights are configurable per job. When graph is sparse, vector score compensates automatically."
    run3.font.size = Pt(10)
    run3.font.italic = True
    run3.font.color.rgb = MED_GRAY
    run3.font.name = "Calibri"

    # Bottom accent line + slide number
    add_shape_rect(slide, Inches(0), Inches(7.25), SLIDE_WIDTH, Inches(0.05), ACCENT_BLUE)
    tb_num = add_text_box(slide, Inches(12.4), Inches(7.05), Inches(0.8), Inches(0.35))
    set_text(tb_num.shape, "2 / 6", font_size=10, color=MED_GRAY, alignment=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════
#  SLIDE 3 — Winning Features #2 (Explainability) & #3 (Bias)
# ══════════════════════════════════════════════════════════════════
def slide_3_features(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # Header
    add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.95), DARK_BG)
    tb = add_text_box(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.6))
    set_text(tb.shape, "Winning Features #2 & #3: Full Explainability + Bias-Aware Screening", font_size=24, bold=True, color=WHITE)

    # ── Left: Explainability ──
    tb_wf2 = add_text_box(slide, Inches(0.4), Inches(1.15), Inches(6), Inches(0.4))
    set_text(tb_wf2.shape, "FEATURE #2: FULL-PIPELINE EXPLAINABILITY", font_size=14, bold=True, color=ACCENT_GREEN)

    explain_data = [
        ["Layer", "What RAX Explains", "Example"],
        ["Skill Matching", "Named graph paths for each skill", "4/5 must-have: Python ✅, SQL ✅, AWS ✅, Docker ✅, K8s ❌"],
        ["Similar Skills", "Adjacent skills via graph + embeddings", "React → 87% partial credit for Frontend"],
        ["Experience Depth", "Years comparison per skill", "Python: has 8, needs 5 → surplus; K8s: has 0, needs 2 → gap"],
        ["Education Fit", "Degree level comparison", "Master's exceeds required Bachelor's level"],
        ["Overall Narrative", "Natural language AI explanation", "Strong match. Gap in K8s. Recommend interview."],
        ["Bias Audit", "Queryable proof of fair scoring", "Graph query: no protected attributes in score path"],
    ]
    add_table(slide, Inches(0.4), Inches(1.55), Inches(6.2), Inches(3.1), 7, 3, explain_data,
              col_widths=[Inches(1.2), Inches(2.2), Inches(2.8)], header_color=ACCENT_GREEN)

    # Deliverable callout
    del_box = add_rounded_rect(slide, Inches(0.4), Inches(4.85), Inches(6.2), Inches(0.75), RGBColor(0xEA, 0xFB, 0xEF), ACCENT_GREEN)
    tf_d = del_box.text_frame
    tf_d.word_wrap = True
    p = tf_d.paragraphs[0]
    run = p.add_run()
    run.text = "Recruiter sees: "
    run.font.bold = True
    run.font.size = Pt(11)
    run.font.color.rgb = ACCENT_GREEN
    run.font.name = "Calibri"
    run2 = p.add_run()
    run2.text = "Radar chart (skills/exp/edu) + Strengths/Gaps cards + Full AI narrative + \"Reveal Identity\" toggle for blind screening"
    run2.font.size = Pt(11)
    run2.font.color.rgb = BODY_TEXT
    run2.font.name = "Calibri"

    # ── Right: Bias-Aware Screening ──
    tb_wf3 = add_text_box(slide, Inches(6.9), Inches(1.15), Inches(6.2), Inches(0.4))
    set_text(tb_wf3.shape, "FEATURE #3: BIAS-AWARE SCREENING BY DESIGN", font_size=14, bold=True, color=PURPLE)

    # Pipeline flow with anonymization
    flow_box = add_rounded_rect(slide, Inches(6.9), Inches(1.6), Inches(6.1), Inches(2.2), RGBColor(0xF5, 0xEE, 0xFC), PURPLE)
    tf_flow = flow_box.text_frame
    tf_flow.word_wrap = True
    set_text(flow_box, "Anonymization Pipeline:", font_size=12, bold=True, color=PURPLE)
    add_paragraph(tf_flow, "", font_size=4)
    add_paragraph(tf_flow, "Original Resume → ResumeParserAgent → Structured JSON", font_size=11, color=DARK_TEXT, font_name="Consolas")
    add_paragraph(tf_flow, "                                           ↓", font_size=11, color=DARK_TEXT, font_name="Consolas")
    add_paragraph(tf_flow, "                                    BiasFilterAgent", font_size=11, bold=True, color=PURPLE, font_name="Consolas")
    add_paragraph(tf_flow, "                                           ↓", font_size=11, color=DARK_TEXT, font_name="Consolas")
    add_paragraph(tf_flow, "                               Anonymized Resume", font_size=11, bold=True, color=DARK_TEXT, font_name="Consolas")

    # Bias details
    bias_items = [
        "Anonymization runs BEFORE any scoring or matching occurs",
        "Name → [CANDIDATE_ID], Harvard → [UNIVERSITY]",
        "Gender, nationality, age signals → removed",
        "\"Reveal Identity\" is a conscious choice AFTER reviewing scores",
        "Graph audit: verify no protected attributes influenced score",
        "Compliant with EEOC, DEI mandates, EU AI Act",
    ]
    tb_bias = add_text_box(slide, Inches(6.9), Inches(4.0), Inches(6.1), Inches(2.5))
    tf_bias = tb_bias.text_frame
    tf_bias.word_wrap = True
    add_bullet_list(tf_bias, bias_items, font_size=12, color=BODY_TEXT, bold_first_part=False)

    # Key stat callout
    stat_box = add_rounded_rect(slide, Inches(6.9), Inches(5.95), Inches(6.1), Inches(0.55), RGBColor(0xFC, 0xEE, 0xF5), PINK)
    tf_stat = stat_box.text_frame
    tf_stat.word_wrap = True
    p = tf_stat.paragraphs[0]
    run = p.add_run()
    run.text = "Studies show identical resumes with 'white-sounding' names receive 50% more callbacks. RAX eliminates this by design."
    run.font.size = Pt(11)
    run.font.italic = True
    run.font.color.rgb = PINK
    run.font.name = "Calibri"

    # Bottom accent line + slide number
    add_shape_rect(slide, Inches(0), Inches(7.25), SLIDE_WIDTH, Inches(0.05), ACCENT_BLUE)
    tb_num = add_text_box(slide, Inches(12.4), Inches(7.05), Inches(0.8), Inches(0.35))
    set_text(tb_num.shape, "3 / 6", font_size=10, color=MED_GRAY, alignment=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════
#  SLIDE 4 — Winning Feature #4 (Real-Time Pipeline) + Workflow
# ══════════════════════════════════════════════════════════════════
def slide_4_workflow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # Header
    add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.95), DARK_BG)
    tb = add_text_box(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.6))
    set_text(tb.shape, "Winning Feature #4: Real-Time Multi-Agent Pipeline + Recruiter Workflow", font_size=24, bold=True, color=WHITE)

    # ── Left: Real-Time Pipeline ──
    tb_wf4 = add_text_box(slide, Inches(0.4), Inches(1.15), Inches(6), Inches(0.4))
    set_text(tb_wf4.shape, "REAL-TIME WEBSOCKET PIPELINE", font_size=14, bold=True, color=ORANGE)

    # Live card mock
    card_box = add_rounded_rect(slide, Inches(0.4), Inches(1.6), Inches(6.2), Inches(2.5), RGBColor(0xFE, 0xF9, 0xEF), ORANGE)
    tf_card = card_box.text_frame
    tf_card.word_wrap = True
    set_text(card_box, "  Resume: john_doe_resume.pdf", font_size=13, bold=True, color=DARK_TEXT, font_name="Consolas")
    add_paragraph(tf_card, "", font_size=4)
    stages = [
        ("  ● Parsing .............. ✅ Complete (2.1s)", ACCENT_GREEN),
        ("  ● Bias Filtering ....... ✅ Complete (1.3s)", ACCENT_GREEN),
        ("  ● Graph Ingestion ...... ✅ Complete (1.8s)  ┐ parallel", ACCENT_GREEN),
        ("  ● Vector Embedding ..... ✅ Complete (0.9s)  ┘", ACCENT_GREEN),
        ("  ● Hybrid Matching ...... ✅ Complete (0.6s)", ACCENT_GREEN),
        ("  ● Scoring .............. ✅ Complete (3.2s)", ACCENT_GREEN),
    ]
    for text, color in stages:
        add_paragraph(tf_card, text, font_size=11, color=color, font_name="Consolas")
    add_paragraph(tf_card, "", font_size=4)
    add_paragraph(tf_card, "  Overall: 87/100   Skills: 91   Exp: 84   Edu: 78", font_size=12, bold=True, color=ACCENT_BLUE, font_name="Consolas")

    # Benefits
    pipeline_benefits = [
        "Recruiters see each stage animate in real time via WebSocket",
        "Multiple resumes process concurrently, each with its own card",
        "Failures are visible immediately — no silent rejections",
        "Graph + Embedding run in parallel, saving ~30% pipeline time",
        "Team members can watch the same pipeline in real time",
    ]
    tb_ben = add_text_box(slide, Inches(0.4), Inches(4.3), Inches(6.2), Inches(2.2))
    tf_ben = tb_ben.text_frame
    tf_ben.word_wrap = True
    add_bullet_list(tf_ben, pipeline_benefits, font_size=12, color=BODY_TEXT, bold_first_part=False)

    # ── Right: Recruiter Workflow ──
    tb_wflow = add_text_box(slide, Inches(6.9), Inches(1.15), Inches(6), Inches(0.4))
    set_text(tb_wflow.shape, "END-TO-END RECRUITER WORKFLOW", font_size=14, bold=True, color=ACCENT_BLUE)

    workflow_steps = [
        ("1", "LOGIN", "Authenticate with role (Recruiter / Hiring Manager)", ACCENT_BLUE),
        ("2", "CREATE JOB", "JD embedded in Qdrant + decomposed into Neo4j graph", ACCENT_BLUE),
        ("3", "UPLOAD RESUMES", "Bulk drag-and-drop (PDF/DOCX)", ACCENT_BLUE),
        ("4", "WATCH LIVE", "WebSocket stage cards per resume", ORANGE),
        ("5", "REVIEW RANKED", "Sorted by hybrid score, filterable by range", ACCENT_GREEN),
        ("6a", "VIEW DETAIL", "Radar chart + strengths/gaps + AI explanation", ACCENT_GREEN),
        ("6b", "REVEAL IDENTITY", "Conscious choice after reviewing scores", PURPLE),
        ("6c", "GENERATE FEEDBACK", "AI-written constructive rejection email", PINK),
        ("7", "SEND FEEDBACK", "Copy to clipboard / Mark as sent", PINK),
    ]

    y_pos = 1.58
    for step_num, step_name, step_desc, color in workflow_steps:
        # Step number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.9), Inches(y_pos), Inches(0.4), Inches(0.4))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        tf_c = circle.text_frame
        tf_c.word_wrap = False
        p = tf_c.paragraphs[0]
        p.text = step_num
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER

        # Step text
        tb_step = add_text_box(slide, Inches(7.4), Inches(y_pos - 0.02), Inches(5.6), Inches(0.42))
        tf_step = tb_step.text_frame
        tf_step.word_wrap = True
        p = tf_step.paragraphs[0]
        run1 = p.add_run()
        run1.text = step_name + "  "
        run1.font.size = Pt(12)
        run1.font.bold = True
        run1.font.color.rgb = color
        run1.font.name = "Calibri"
        run2 = p.add_run()
        run2.text = step_desc
        run2.font.size = Pt(11)
        run2.font.color.rgb = BODY_TEXT
        run2.font.name = "Calibri"

        y_pos += 0.57

    # Bottom accent line + slide number
    add_shape_rect(slide, Inches(0), Inches(7.25), SLIDE_WIDTH, Inches(0.05), ACCENT_BLUE)
    tb_num = add_text_box(slide, Inches(12.4), Inches(7.05), Inches(0.8), Inches(0.35))
    set_text(tb_num.shape, "4 / 6", font_size=10, color=MED_GRAY, alignment=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════
#  SLIDE 5 — System Architecture
# ══════════════════════════════════════════════════════════════════
def slide_5_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # Header
    add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.95), DARK_BG)
    tb = add_text_box(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.6))
    set_text(tb.shape, "System Architecture — Four-Tier Cloud-Native Design", font_size=24, bold=True, color=WHITE)

    # ── CLIENT TIER box ──
    client_box = add_rounded_rect(slide, Inches(3.5), Inches(1.15), Inches(6.3), Inches(0.85), RGBColor(0xE3, 0xF2, 0xFD), ACCENT_BLUE)
    tf_cl = client_box.text_frame
    tf_cl.word_wrap = True
    set_text(client_box, "CLIENT TIER  —  React + TypeScript (Vite) + Tailwind + shadcn/ui", font_size=12, bold=True, color=ACCENT_BLUE)
    add_paragraph(tf_cl, "Zustand State  |  Axios REST  |  WebSocket Live Updates  |  Deployed on Vercel", font_size=10, color=BODY_TEXT)

    # Down arrow
    arrow1 = add_text_box(slide, Inches(6.2), Inches(2.0), Inches(1), Inches(0.35))
    set_text(arrow1.shape, "REST + WS  ↕", font_size=9, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

    # ── SERVER TIER box ──
    server_box = add_rounded_rect(slide, Inches(1.8), Inches(2.35), Inches(9.7), Inches(2.3), RGBColor(0xE8, 0xF5, 0xE9), ACCENT_GREEN)
    tf_sv = server_box.text_frame
    tf_sv.word_wrap = True
    set_text(server_box, "SERVER TIER  —  Python FastAPI (Async)  |  JWT Auth  |  Alembic  |  Deployed on Railway / Render", font_size=12, bold=True, color=ACCENT_GREEN)

    # API box inside server
    api_box = add_rounded_rect(slide, Inches(2.0), Inches(2.85), Inches(3.0), Inches(0.65), RGBColor(0xC8, 0xE6, 0xC9), ACCENT_GREEN)
    set_text(api_box, "  Auth Routes  |  Jobs API  |  Resumes API  |  WebSocket Endpoint", font_size=10, color=DARK_TEXT)

    # Pipeline box inside server
    pipe_box = add_rounded_rect(slide, Inches(2.0), Inches(3.65), Inches(9.3), Inches(0.85), RGBColor(0xC8, 0xE6, 0xC9), ACCENT_GREEN)
    tf_pipe = pipe_box.text_frame
    tf_pipe.word_wrap = True
    set_text(pipe_box, "  Pipeline Orchestrator:", font_size=10, bold=True, color=ACCENT_GREEN)
    add_paragraph(tf_pipe, "  ResumeParser → BiasFilter → [ GraphIngestion ‖ Embedding ] → HybridMatch → Scoring   (+FeedbackAgent on demand)", font_size=10, color=DARK_TEXT, font_name="Consolas")

    # ── Bottom Tier Boxes ──
    # Down arrows
    arrow_data = add_text_box(slide, Inches(2.5), Inches(4.75), Inches(1.2), Inches(0.3))
    set_text(arrow_data.shape, "SQL ↕", font_size=9, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
    arrow_graph = add_text_box(slide, Inches(5.2), Inches(4.75), Inches(1.2), Inches(0.3))
    set_text(arrow_graph.shape, "Bolt ↕", font_size=9, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
    arrow_vector = add_text_box(slide, Inches(7.9), Inches(4.75), Inches(1.2), Inches(0.3))
    set_text(arrow_vector.shape, "REST ↕", font_size=9, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
    arrow_ai = add_text_box(slide, Inches(10.6), Inches(4.75), Inches(1.2), Inches(0.3))
    set_text(arrow_ai.shape, "HTTPS ↕", font_size=9, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

    # DATA TIER
    data_box = add_rounded_rect(slide, Inches(0.4), Inches(5.1), Inches(3.0), Inches(1.5), RGBColor(0xFF, 0xF3, 0xE0), ORANGE)
    tf_data = data_box.text_frame
    tf_data.word_wrap = True
    set_text(data_box, "  DATA TIER", font_size=11, bold=True, color=ORANGE)
    add_paragraph(tf_data, "  Supabase (PostgreSQL)", font_size=10, bold=True, color=DARK_TEXT)
    add_paragraph(tf_data, "  Users | Jobs | Resumes", font_size=9, color=BODY_TEXT)
    add_paragraph(tf_data, "  Analyses | Feedback", font_size=9, color=BODY_TEXT)
    add_paragraph(tf_data, "  + Supabase Storage (files)", font_size=9, color=BODY_TEXT)

    # GRAPH TIER
    graph_box = add_rounded_rect(slide, Inches(3.7), Inches(5.1), Inches(3.0), Inches(1.5), RGBColor(0xF3, 0xE5, 0xF5), PURPLE)
    tf_graph = graph_box.text_frame
    tf_graph.word_wrap = True
    set_text(graph_box, "  GRAPH TIER", font_size=11, bold=True, color=PURPLE)
    add_paragraph(tf_graph, "  Neo4j (AuraDB)", font_size=10, bold=True, color=DARK_TEXT)
    add_paragraph(tf_graph, "  Candidate | Skill | Company", font_size=9, color=BODY_TEXT)
    add_paragraph(tf_graph, "  Role | Education | SkillCluster", font_size=9, color=BODY_TEXT)
    add_paragraph(tf_graph, "  HAS_SKILL | IS_SIMILAR_TO", font_size=9, color=BODY_TEXT)

    # VECTOR TIER
    vector_box = add_rounded_rect(slide, Inches(7.0), Inches(5.1), Inches(3.0), Inches(1.5), RGBColor(0xFC, 0xE4, 0xEC), PINK)
    tf_vec = vector_box.text_frame
    tf_vec.word_wrap = True
    set_text(vector_box, "  VECTOR TIER", font_size=11, bold=True, color=PINK)
    add_paragraph(tf_vec, "  Qdrant (Vector DB)", font_size=10, bold=True, color=DARK_TEXT)
    add_paragraph(tf_vec, "  resumes collection", font_size=9, color=BODY_TEXT)
    add_paragraph(tf_vec, "  job_descs collection", font_size=9, color=BODY_TEXT)
    add_paragraph(tf_vec, "  Cosine similarity search", font_size=9, color=BODY_TEXT)

    # AI TIER
    ai_box = add_rounded_rect(slide, Inches(10.3), Inches(5.1), Inches(2.7), Inches(1.5), RGBColor(0xFD, 0xED, 0xED), ACCENT_RED)
    tf_ai = ai_box.text_frame
    tf_ai.word_wrap = True
    set_text(ai_box, "  AI TIER", font_size=11, bold=True, color=ACCENT_RED)
    add_paragraph(tf_ai, "  Google Gemini API", font_size=10, bold=True, color=DARK_TEXT)
    add_paragraph(tf_ai, "  gemini-1.5-pro", font_size=9, color=BODY_TEXT)
    add_paragraph(tf_ai, "  (parse, score, feedback)", font_size=9, color=BODY_TEXT)
    add_paragraph(tf_ai, "  text-embedding-004", font_size=9, color=BODY_TEXT)

    # Distributed computing note
    dist_box = add_rounded_rect(slide, Inches(5.2), Inches(2.85), Inches(6.1), Inches(0.65), RGBColor(0xFE, 0xF9, 0xEF), ORANGE)
    tf_dist = dist_box.text_frame
    tf_dist.word_wrap = True
    set_text(dist_box, "  Distributed: Parallel agents | Async pipeline | Stateless JWT | Self-enriching graph | WebSocket pub/sub | Cloud-native ($0)", font_size=9, color=ORANGE)

    # Bottom accent line + slide number
    add_shape_rect(slide, Inches(0), Inches(7.25), SLIDE_WIDTH, Inches(0.05), ACCENT_BLUE)
    tb_num = add_text_box(slide, Inches(12.4), Inches(7.05), Inches(0.8), Inches(0.35))
    set_text(tb_num.shape, "5 / 6", font_size=10, color=MED_GRAY, alignment=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════
#  SLIDE 6 — Conclusion & Implementation Plan
# ══════════════════════════════════════════════════════════════════
def slide_6_conclusion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # Header
    add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.95), DARK_BG)
    tb = add_text_box(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.6))
    set_text(tb.shape, "Conclusion & Implementation Plan", font_size=24, bold=True, color=WHITE)

    # ── Left: Why RAX Wins ──
    tb_label = add_text_box(slide, Inches(0.4), Inches(1.15), Inches(6), Inches(0.4))
    set_text(tb_label.shape, "WHY RAX WINS", font_size=14, bold=True, color=ACCENT_BLUE)

    wins_data = [
        ["What Competitors Do", "What RAX Does Differently"],
        ["Keyword matching", "Semantic understanding (Gemini LLM)"],
        ["Black-box scoring", "Full explainability (graph paths + AI narrative)"],
        ["Bias as afterthought", "Anonymization as core pipeline stage"],
        ["Vector-only similarity", "Hybrid graph + vector matching"],
        ["Batch processing", "Real-time WebSocket with live stage visibility"],
        ["No candidate feedback", "AI-generated constructive rejection feedback"],
        ["Expensive enterprise tools", "100% free-tier cloud deployment ($0)"],
    ]
    add_table(slide, Inches(0.4), Inches(1.55), Inches(6.2), Inches(3.2), 8, 2, wins_data,
              col_widths=[Inches(2.5), Inches(3.7)], header_color=ACCENT_BLUE)

    # Risk mitigation
    tb_risk = add_text_box(slide, Inches(0.4), Inches(4.95), Inches(6.2), Inches(0.4))
    set_text(tb_risk.shape, "RISK MITIGATION", font_size=12, bold=True, color=ACCENT_RED)

    risk_items = [
        "Gemini rate limits → retry + exponential backoff + cache embeddings",
        "Neo4j 200K node limit → normalize/dedup skill names",
        "Sparse graph data → Qdrant vector score auto-compensates",
        "WebSocket drops → client-side reconnect + polling fallback",
    ]
    tb_risks = add_text_box(slide, Inches(0.4), Inches(5.35), Inches(6.2), Inches(1.7))
    tf_risks = tb_risks.text_frame
    tf_risks.word_wrap = True
    add_bullet_list(tf_risks, risk_items, font_size=11, color=BODY_TEXT, bold_first_part=False)

    # ── Right: Implementation Plan ──
    tb_plan = add_text_box(slide, Inches(6.9), Inches(1.15), Inches(6), Inches(0.4))
    set_text(tb_plan.shape, "IMPLEMENTATION PLAN (6 WEEKS)", font_size=14, bold=True, color=ACCENT_GREEN)

    plan_data = [
        ["Phase", "Scope", "Who", "When"],
        ["0", "Repo scaffold, Docker (Neo4j + Qdrant), Supabase", "All", "Week 1"],
        ["1", "Backend: FastAPI, DB, Auth, CRUD, Clients", "Person 1", "Weeks 2–3"],
        ["2", "AI Pipeline: all 6 agents + orchestrator", "Person 2", "Weeks 2–4"],
        ["3", "WebSocket real-time broadcasting", "Person 2", "Week 4"],
        ["4", "Frontend: all screens + WS integration", "Person 3", "Weeks 2–5"],
        ["5", "Cloud deploy (Vercel, Railway, AuraDB, Qdrant)", "All", "Week 5"],
        ["6", "E2E integration testing + demo prep", "All", "Week 6"],
    ]
    add_table(slide, Inches(6.9), Inches(1.55), Inches(6.1), Inches(2.8), 8, 4, plan_data,
              col_widths=[Inches(0.6), Inches(3.0), Inches(1.0), Inches(1.5)], header_color=ACCENT_GREEN)

    # Milestones
    tb_ms = add_text_box(slide, Inches(6.9), Inches(4.55), Inches(6.1), Inches(0.35))
    set_text(tb_ms.shape, "KEY MILESTONES", font_size=12, bold=True, color=ORANGE)

    milestones = [
        "W1: Infra ready (Docker + Supabase + repo scaffolded)",
        "W3: Backend API live (auth + CRUD + databases connected)",
        "W4: Full AI pipeline works (upload → hybrid scores + explanation)",
        "W5: Complete recruiter workflow in frontend",
        "W6: Deployed & verified end-to-end on cloud infrastructure",
    ]
    tb_mile = add_text_box(slide, Inches(6.9), Inches(4.9), Inches(6.1), Inches(1.8))
    tf_mile = tb_mile.text_frame
    tf_mile.word_wrap = True
    add_bullet_list(tf_mile, milestones, font_size=11, color=BODY_TEXT, bold_first_part=False)

    # Closing statement box
    close_box = add_rounded_rect(slide, Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.55), DARK_BG)
    tf_close = close_box.text_frame
    tf_close.word_wrap = True
    p = tf_close.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "RAX: The first hiring platform combining structural reasoning + semantic discovery. Auditable. Bias-free. Real-time. Ready to build."
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = ACCENT_TEAL
    run.font.name = "Calibri"

    # Bottom accent line + slide number
    add_shape_rect(slide, Inches(0), Inches(7.25), SLIDE_WIDTH, Inches(0.05), ACCENT_BLUE)
    tb_num = add_text_box(slide, Inches(12.4), Inches(7.05), Inches(0.8), Inches(0.35))
    set_text(tb_num.shape, "6 / 6", font_size=10, color=MED_GRAY, alignment=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide_1_title(prs)
    slide_2_competitive(prs)
    slide_3_features(prs)
    slide_4_workflow(prs)
    slide_5_architecture(prs)
    slide_6_conclusion(prs)

    output_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(output_dir, "RAX_Stage2_Presentation.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")


if __name__ == "__main__":
    main()
